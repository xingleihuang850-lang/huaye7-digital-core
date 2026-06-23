#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""花页7 PPT 汇报 → deliverables/花页7/花页7_多尺度数字岩心汇报.pptx"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT=os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ST=json.load(open(os.path.join(ROOT,"experiments","hy7_stats.json"),encoding="utf-8"))
FIG=os.path.join(ROOT,"experiments","hy7_figures")
OUT=os.path.join(ROOT,"deliverables","花页7","花页7_多尺度数字岩心汇报.pptx")
NAVY=RGBColor(0x10,0x23,0x3A); DEEP=RGBColor(0x0E,0x4A,0x66); TEAL=RGBColor(0x1C,0x72,0x93)
AMBER=RGBColor(0xE8,0xA3,0x3D); ICE=RGBColor(0xCB,0xDD,0xEA); WHITE=RGBColor(0xFF,0xFF,0xFF)
INK=RGBColor(0x22,0x2B,0x33); MUTE=RGBColor(0x6B,0x7B,0x88); LIGHT=RGBColor(0xF3,0xF7,0xFA); CARD=RGBColor(0xFF,0xFF,0xFF)
FONT="Arial"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def rect(s,x,y,w,h,color):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background(); r.shadow.inherit=False; return r
def rrect(s,x,y,w,h,color):
    r=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background(); r.shadow.inherit=False; return r
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(2); tf.margin_top=tf.margin_bottom=Pt(2)
    if isinstance(runs,str): runs=[(runs,16,INK,False)]
    first=True
    for text,size,color,bold in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.line_spacing=space; p.space_after=Pt(2)
        r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=FONT
    return tb
def pic(s,fn,x,y,mw,mh):
    p=os.path.join(FIG,fn); im=Image.open(p); a=im.size[0]/im.size[1]
    w=mw; h=w/a
    if h>mh: h=mh; w=h*a
    x=x+(mw-w)/2; y=y+(mh-h)/2
    s.shapes.add_picture(p,Inches(x),Inches(y),Inches(w),Inches(h))
def card(s,x,y,w,h,num,lab):
    rrect(s,x,y,w,h,CARD)
    txt(s,x,y+0.12,w,h*0.55,[(num,26,DEEP,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.06,y+h*0.6,w-0.12,h*0.36,[(lab,10.5,MUTE,False)],align=PP_ALIGN.CENTER)

g=ST["lithology_groups"]
# 1 标题
s=slide(NAVY); rect(s,0,0,13.333,0.18,AMBER)
txt(s,0.9,2.2,11.5,1.3,[("花页7井 多尺度数字岩心",46,WHITE,True)])
txt(s,0.92,3.5,11.5,0.7,[("6 个成像尺度，从 25μm 矿物到 8.589nm 纳米孔",22,ICE,False)])
txt(s,0.92,4.25,11.5,0.6,[("HuaYe-7 Multi-scale Digital Rock · 4199.21 m 页岩",14,AMBER,False)])
txt(s,0.92,6.5,11.5,0.5,[("数据已核校（12 项全过）· 多尺度对齐 + 矿物分割 · 2026-06-22",13,ICE,False)])
# 2 速览
s=slide(LIGHT)
txt(s,0.7,0.45,12,0.8,[("一页速览：这是什么数据",34,NAVY,True)])
txt(s,0.72,1.32,12,0.5,[("吉林大学花页7井 4199.21m 页岩，6 个成像尺度全做完三维孔隙/喉道/裂缝分割与矿物定量。",15,INK,False)])
cards=[("4199.21 m","深度·页岩"),("6 个","成像尺度"),("25μm→8.6nm","跨~3000倍"),
       ("7.18%","最高总孔隙度(2.8μm)"),("76,252","FIB-SEM 孔隙数"),("12/12","核校通过")]
for i,(n,l) in enumerate(cards): card(s,0.7+i*2.08,2.25,1.95,1.5,n,l)
txt(s,0.7,4.1,12,0.5,[("三句话记住它",20,TEAL,True)])
txt(s,0.72,4.7,12,2.2,[
    ("①  六个尺度看同一块页岩的不同孔隙层级，需多尺度对齐+融合（非相加）。",16,INK,False),
    ("②  总孔隙度非单调（2.8μm 最高 7.18%）= 尺度效应/REV，是讨论点不是错误。",16,INK,False),
    (f"③  钙质长英质页岩：长英质{g['长英质']}%+碳酸盐{g['碳酸盐']}%+黏土{g['黏土']}%；黏土低。",16,INK,False),
],space=1.25)
# 3 尺度阶梯
s=slide(WHITE)
txt(s,0.7,0.45,12,0.8,[("六尺度阶梯：25μm → 8.589nm",32,NAVY,True)])
pic(s,"06_ladder.png",0.5,1.5,8.4,5.3)
txt(s,9.1,1.6,3.9,0.6,[("怎么理解",20,TEAL,True)])
txt(s,9.12,2.25,4.0,4.6,[
    ("每个尺度看一层孔隙系统：",14,INK,False),
    ("",6,INK,False),
    ("· Amics 25μm → 矿物组成",13,MUTE,False),
    ("· 微米CT 14/2.8μm → 骨架孔、裂缝",13,MUTE,False),
    ("· 纳米CT 65nm → 纳米孔",13,MUTE,False),
    ("· Maps 10nm → 有机/无机孔",13,MUTE,False),
    ("· FIB-SEM 8.6nm → 孔喉网络",13,MUTE,False),
    ("",6,INK,False),
    ("跨约 3000 倍，拼起来才是完整数字岩心。",14,INK,False),
],space=1.12)
# 4 多尺度孔隙度
s=slide(LIGHT)
txt(s,0.7,0.45,12,0.8,[("多尺度总孔隙度：为什么非单调",32,NAVY,True)])
pic(s,"01_porosity.png",0.5,1.45,8.5,5.2)
rrect(s,9.35,1.55,3.5,5.3,CARD)
txt(s,9.55,1.75,3.1,5.0,[
    ("尺度效应 / REV",18,AMBER,True),("",6,INK,False),
    ("2.8μm 最高 7.18%——恰好捕捉微孔-微裂缝主峰。",14,INK,False),("",5,INK,False),
    ("不同尺度视场、可分辨孔径、统计体不同，总孔隙度不可简单相加。",13,MUTE,False),("",5,INK,False),
    ("注：FIB-SEM 报的 1.52% 是有机质占比、非孔隙度，故不计入对比。",12,MUTE,False),("",5,INK,False),
    ("→ 这正是要做多尺度融合的证据。",14,TEAL,True),
],space=1.1)
# 5 矿物
s=slide(WHITE)
txt(s,0.7,0.45,12,0.8,[("矿物组成：钙质长英质页岩",32,NAVY,True)])
pic(s,"02_minerals.png",0.4,1.5,8.7,4.6)
txt(s,9.2,1.6,3.8,0.5,[("要点",20,TEAL,True)])
txt(s,9.22,2.2,3.9,4.6,[
    (f"长英质 ≈{g['长英质']}%",16,DEEP,True),
    ("（斜长石31+石英24+钾长石）",12,MUTE,False),("",4,INK,False),
    (f"碳酸盐 ≈{g['碳酸盐']}%",16,AMBER,True),
    ("（方解石18+白云石12+铁白云石11）",12,MUTE,False),("",4,INK,False),
    (f"黏土 仅 ≈{g['黏土']}%",16,DEEP,True),
    ("黄铁矿 0.48%",12,MUTE,False),("",6,INK,False),
    ("⚠ 粗扫(25μm)与精扫(1μm)比例不同，",13,INK,False),
    ("矿物分割要先定标定基准。",13,INK,False),
],space=1.12)
# 6 孔喉结构
s=slide(LIGHT)
txt(s,0.7,0.45,12,0.8,[("孔喉结构与连通性",32,NAVY,True)])
pic(s,"03_pore_radius.png",0.4,1.45,7.4,5.5)
pic(s,"04_coordination.png",7.9,1.5,5.1,3.4)
txt(s,8.0,5.1,5.0,1.9,[
    ("连通性整体偏弱：配位数集中在低值、0 为最大占比（14μm 达 85%，其余约 36–47%），孤立孔较多。",14,INK,False),
    (f"FIB-SEM 网络：孔隙 {ST['fib_counts']['pores']:,} 个、喉道 {ST['fib_counts']['throats']:,} 个。",14,DEEP,True),
],space=1.12)
# 7 任务
s=slide(NAVY); rect(s,0,0,13.333,0.18,AMBER)
txt(s,0.7,0.5,12,0.8,[("核心任务：多尺度对齐 + 矿物分割",32,WHITE,True)])
txt(s,0.72,1.4,12,0.5,[("数据现成、判别式分割已完成；下一跳是把尺度对齐、矿物分好、再生成式融合。",15,ICE,False)])
tasks=[("多尺度对齐(配准)","真瓶颈：图像-标签未配准","Plan B：用天然对齐的 sus.raw+pore.raw 训练，绕开原图配准"),
       ("矿物分割","已有 Amics 矿物图，粗/精扫基准不一","定一个标定基准，迁移到 CT 体做长英质/碳酸盐/黏土分类"),
       ("多尺度融合+生成","各尺度总孔隙度非单调、各看一层","跨尺度孔网拼接 + 体素超分辨/生成式补全 + REV")]
for i,(t,now,how) in enumerate(tasks):
    cx=0.7+i*4.17; rrect(s,cx,2.2,3.9,4.4,RGBColor(0x18,0x33,0x50))
    txt(s,cx+0.25,2.45,3.4,0.9,[(f"任务 {i+1}",13,AMBER,True),(t,17,WHITE,True)])
    txt(s,cx+0.25,3.8,3.4,1.3,[("现状",12,MUTE,True),(now,13,ICE,False)],space=1.05)
    txt(s,cx+0.25,5.25,3.4,1.2,[("怎么做",12,AMBER,True),(how,13,WHITE,False)],space=1.05)
# 8 闭幕
s=slide(NAVY); rect(s,0,0,13.333,0.18,AMBER)
txt(s,0.9,1.5,11.5,1.0,[("一句话总结",24,AMBER,True)])
txt(s,0.9,2.4,11.6,2.6,[
    ("花页7 是一套现成的、判别式分好的多尺度数字岩心——",26,WHITE,True),
    ("6 个尺度、矿物齐全、孔喉网络清晰，且已逐项核校可信。",26,WHITE,True),
    ("",12,WHITE,False),
    ("缺的是把尺度对齐、把矿物分好、再生成式融合：",24,ICE,False),
    ("从几个尺度的切片，拼成一整块可外推的数字岩心。",24,ICE,False),
],space=1.2)
txt(s,0.92,6.6,11.5,0.5,[("花页7井 4199.21m · 多尺度数字岩心 · 生成式数字井筒研究",13,MUTE,False)])

os.makedirs(os.path.dirname(OUT),exist_ok=True); prs.save(OUT)
print(">> 已生成",os.path.relpath(OUT,ROOT),"  共",len(prs.slides._sldIdLst),"页")
