#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 PPT 汇报演示：deliverables/GJ5-15_数字井筒汇报.pptx
数据来源：experiments/master_stats.json + experiments/figures/*.png
用法： .venv/bin/python src/make_ppt.py
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ST = json.load(open(os.path.join(ROOT, "experiments", "master_stats.json"), encoding="utf-8"))
FIG = os.path.join(ROOT, "experiments", "figures")
OUT = os.path.join(ROOT, "deliverables", "GJ5-15_数字井筒汇报.pptx")
META = ST["curve_meta"]; SEGS = ST["segments"]; CURVE_ORDER = ST["curve_order"]; PS = ST["per_segment"]

# 配色：地下/流体主题 —— 深海军蓝 + 青 + 琥珀(油)
NAVY = RGBColor(0x10, 0x23, 0x3A)
DEEP = RGBColor(0x0E, 0x4A, 0x66)
TEAL = RGBColor(0x1C, 0x72, 0x93)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
ICE = RGBColor(0xCB, 0xDD, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x2B, 0x33)
MUTE = RGBColor(0x6B, 0x7B, 0x88)
LIGHTBG = RGBColor(0xF3, 0xF7, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def rect(s, x, y, w, h, color, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line: r.line.color.rgb = line; r.line.width = Pt(1)
    else: r.line.fill.background()
    r.shadow.inherit = False
    return r


def rrect(s, x, y, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    r.shadow.inherit = False
    return r


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str): runs = [(runs, 16, INK, False)]
    first = True
    for item in runs:
        text, size, color, bold = (item + (False,))[:4] if len(item) < 4 else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = space; p.space_after = Pt(2)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return tb


def pic_fit(s, fname, x, y, max_w, max_h, align="center"):
    """按比例缩放图片放入 (max_w,max_h) 框，可选居中。返回放置矩形。"""
    p = os.path.join(FIG, fname)
    im = Image.open(p); aspect = im.size[0] / im.size[1]
    w = max_w; h = w / aspect
    if h > max_h:
        h = max_h; w = h * aspect
    if align == "center":
        x = x + (max_w - w) / 2; y = y + (max_h - h) / 2
    s.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    return x, y, w, h


def stat_card(s, x, y, w, h, number, label, num_color=DEEP):
    rrect(s, x, y, w, h, CARD)
    txt(s, x, y + 0.12, w, h * 0.55, [(number, 30, num_color, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.08, y + h * 0.6, w - 0.16, h * 0.36, [(label, 11, MUTE, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# ============ 幻灯片 1：标题（深色） ============
s = slide(NAVY)
rect(s, 0, 0, 13.333, 0.18, AMBER)
txt(s, 0.9, 2.25, 11.5, 1.4, [("GJ5-15 数字井筒", 50, WHITE, True)])
txt(s, 0.92, 3.5, 11.5, 0.7, [("从判别式岩心分割，到生成式井筒重建", 22, ICE, False)])
txt(s, 0.92, 4.25, 11.5, 0.6, [("Generative Digital Wellbore — CT 反演数据集解读", 14, AMBER, False)])
txt(s, 0.92, 6.5, 11.5, 0.5,
    [(f"井段 {ST['interval_m'][0]}–{ST['interval_m'][1]} m   ·   5 个深度段   ·   9 条连续曲线   ·   2026-06-22", 13, ICE, False)])

# ============ 幻灯片 2：一页速览 ============
s = slide(LIGHTBG)
txt(s, 0.7, 0.45, 12, 0.8, [("一页速览：这是什么数据", 34, NAVY, True)])
txt(s, 0.72, 1.32, 12, 0.5,
    [("吉林大学 GJ5-15 井 4.4 m 岩心，经 CT 扫描 + 相分割后反演，得到沿井深近连续的多参数“数字井筒”。", 15, INK, False)])
cards = [("4.39 m", "岩心总井段长度"), ("5 段", "深度分段统计"),
         ("9 条", "连续测井曲线"), ("17,982", "每条曲线采样点"),
         ("0.24 mm", "垂向分辨率"), ("50.62 μm", "CT 像素尺寸")]
x0 = 0.7; cw = 1.95; gap = 0.13
for i, (num, lab) in enumerate(cards):
    stat_card(s, x0 + i * (cw + gap), 2.25, cw, 1.5, num, lab)
txt(s, 0.7, 4.1, 12, 0.5, [("三句话记住它", 20, TEAL, True)])
txt(s, 0.72, 4.7, 12, 2.2, [
    ("①  甜点在第 4 段（3441.70–3442.64 m）——含油、裂缝、渗透率全井最高。", 16, INK, False),
    ("②  泥质含量自上而下递减，物性与含油性随之变好。", 16, INK, False),
    ("③  现状是判别式分割成品；研究贡献是“生成式”补全、超分辨与不确定性量化。", 16, INK, False),
], space=1.25)

# ============ 幻灯片 3：综合测井道 ============
s = slide(WHITE)
txt(s, 0.7, 0.45, 12, 0.8, [("数字井筒：九条连续曲线沿井深", 32, NAVY, True)])
pic_fit(s, "01_well_tracks.png", 0.5, 1.35, 7.6, 5.8, align="center")
txt(s, 8.4, 1.5, 4.4, 0.6, [("怎么读这张图", 20, TEAL, True)])
txt(s, 8.42, 2.15, 4.5, 4.8, [
    ("纵轴是深度（向下加深），每一道是一条曲线。", 14, INK, False),
    ("", 6, INK, False),
    ("浅色 = 逐层原始值（约 0.24 mm 一个点）", 13, MUTE, False),
    ("深色 = 平滑趋势线", 13, MUTE, False),
    ("虚线 = 5 个深度段的分界", 13, MUTE, False),
    ("", 6, INK, False),
    ("这就是“数字井筒”的核心：把几块离散岩心，", 14, INK, False),
    ("变成了一条可连续解释的虚拟测井。", 14, INK, False),
], space=1.15)

# ============ 幻灯片 4：分段对比 + 甜点 ============
s = slide(LIGHTBG)
txt(s, 0.7, 0.45, 12, 0.8, [("分段对比：第 4 段是储集甜点", 32, NAVY, True)])
pic_fit(s, "02_segment_bars.png", 0.5, 1.35, 8.5, 5.7, align="center")
seg4 = SEGS[3]
# 右侧甜点卡片
rrect(s, 9.35, 1.5, 3.5, 5.2, CARD)
txt(s, 9.55, 1.7, 3.1, 0.6, [("甜点段", 18, AMBER, True), ("3441.70–3442.64 m", 13, INK, True)])
metrics = [("含油率", PS[seg4]['oil_content']['mean'], "%"),
           ("含油饱和度", PS[seg4]['oil_saturation']['mean'], "%"),
           ("裂缝孔隙度", PS[seg4]['fracture_porosity']['mean'], "%"),
           ("泥质含量(最低)", PS[seg4]['shale_content']['mean'], "%"),
           ("水平渗透率 Kh", 87.8, " mD")]
yy = 2.75
for name, val, unit in metrics:
    txt(s, 9.55, yy, 3.1, 0.32, [(name, 12, MUTE, False)])
    txt(s, 9.55, yy + 0.28, 3.1, 0.45, [(f"{val:.2f}{unit}", 20, DEEP, True)])
    yy += 0.78

# ============ 幻灯片 5：孔隙结构 + 相关性 ============
s = slide(WHITE)
txt(s, 0.7, 0.45, 12, 0.8, [("孔隙结构与参数关系", 32, NAVY, True)])
pic_fit(s, "03_pore_diameter.png", 0.4, 1.5, 6.3, 3.5, align="center")
pic_fit(s, "04_correlation.png", 0.4, 5.0, 6.3, 2.3, align="center")
txt(s, 7.1, 1.5, 5.7, 0.5, [("关键结论", 20, TEAL, True)])
txt(s, 7.12, 2.1, 5.8, 5.0, [
    ("孔隙：数量上小孔（50–100μm）占约 92%，", 15, INK, False),
    ("但大孔贡献了约 40% 的孔隙体积——", 15, INK, False),
    ("储集空间由少数大孔主导。", 15, AMBER, True),
    ("", 8, INK, False),
    ("相关性（逐层 Pearson）：", 15, INK, True),
    (f"· 含油率 ↔ 含油饱和度  r = {ST['correlation']['oil_content']['oil_saturation']:.2f}", 14, INK, False),
    (f"· 泥质 ↔ 含油饱和度  r = {ST['correlation']['shale_content']['oil_saturation']:.2f}", 14, INK, False),
    (f"· 基质孔隙度 ↔ 总孔隙度  r = {ST['correlation']['matrix_porosity']['total_porosity']:.2f}", 14, INK, False),
    ("", 8, INK, False),
    ("→ 泥质是含油性的负向控制因素；", 14, MUTE, False),
    ("  建模时可作为天然的条件变量。", 14, MUTE, False),
], space=1.1)

# ============ 幻灯片 6：渗透率 + 储层评价 ============
s = slide(LIGHTBG)
txt(s, 0.7, 0.45, 12, 0.8, [("渗透率剖面与储层评价", 32, NAVY, True)])
pic_fit(s, "05_permeability.png", 0.4, 1.45, 6.0, 5.6, align="center")
pic_fit(s, "06_reservoir_grade.png", 6.7, 1.45, 4.0, 5.6, align="center")
txt(s, 10.8, 1.6, 2.4, 0.5, [("要点", 18, TEAL, True)])
txt(s, 10.82, 2.15, 2.45, 5.0, [
    ("Kh 在第 4 段冲到 87.8 mD 峰值后回落。", 13, INK, False),
    ("", 6, INK, False),
    ("Kh/Kv 各向异性同段最高（6.59），", 13, INK, False),
    ("层理/裂缝定向性强。", 13, MUTE, False),
    ("", 6, INK, False),
    ("储层评价 a→d 与含油饱和度趋势一致。", 13, INK, False),
], space=1.1)

# ============ 幻灯片 7：判别式 → 生成式（深色） ============
s = slide(NAVY)
rect(s, 0, 0, 13.333, 0.18, AMBER)
txt(s, 0.7, 0.5, 12, 0.8, [("下一跳：从判别式分割到生成式重建", 32, WHITE, True)])
txt(s, 0.72, 1.4, 12, 0.5, [("数据本身暴露了三个缺口，正是生成式建模的发力点。", 15, ICE, False)])
gaps = [
    ("后 2 段无 2D 切片", "3441.70–3443.52 m 仅有统计与 3D 渲染", "用前 3 段切片训练，沿井深生成缺失段虚拟岩心"),
    ("无矿物种类标注", "只有相分割，未区分石英/长石/方解石/黏土", "2D(SEM-EDS) → 3D(CT) 迁移学习生成矿物分布"),
    ("分辨率有限 / 单点确定", "50 μm 难分辨细结构；每深度仅一个值", "体素超分辨 + 多样本生成做不确定性量化"),
]
cw = 3.95; x = 0.7
for i, (title, now, gen) in enumerate(gaps):
    cx = x + i * (cw + 0.22)
    rrect(s, cx, 2.2, cw, 4.4, RGBColor(0x18, 0x33, 0x50))
    txt(s, cx + 0.25, 2.45, cw - 0.5, 0.9, [(f"缺口 {i+1}", 13, AMBER, True), (title, 17, WHITE, True)])
    txt(s, cx + 0.25, 3.75, cw - 0.5, 1.3, [("现状", 12, MUTE, True), (now, 13, ICE, False)], space=1.05)
    txt(s, cx + 0.25, 5.25, cw - 0.5, 1.2, [("生成式怎么补", 12, AMBER, True), (gen, 13, WHITE, False)], space=1.05)

# ============ 幻灯片 8：研究路线 ============
s = slide(WHITE)
txt(s, 0.7, 0.45, 12, 0.8, [("研究路线与下一步", 32, NAVY, True)])
steps = [
    ("1", "建库与基线", "DuckDB 仓库已建（curves / segment_stats 等）；曲线补全先做线性插值 baseline"),
    ("2", "数据补充", "补 SEM-EDS / XRD 提供矿物 Ground Truth；争取后两段原始 CT 切片"),
    ("3", "生成式建模", "体素超分辨（VAE/扩散）、坏井段补全、沿井深虚拟岩心生成（迁移学习最匹配）"),
    ("4", "不确定性量化", "多样本生成给分布而非单点，用前 3 段→后 2 段做天然留出验证"),
]
yy = 1.6
for n, title, desc in steps:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(yy), Inches(0.85), Inches(0.85))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL; circ.line.fill.background(); circ.shadow.inherit = False
    txt(s, 0.8, yy, 0.85, 0.85, [(n, 26, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, 1.9, yy - 0.05, 10.9, 1.05, LIGHTBG)
    txt(s, 2.15, yy + 0.05, 10.5, 0.45, [(title, 18, NAVY, True)])
    txt(s, 2.15, yy + 0.52, 10.5, 0.45, [(desc, 13.5, INK, False)])
    yy += 1.32

# ============ 幻灯片 9：闭幕（深色） ============
s = slide(NAVY)
rect(s, 0, 0, 13.333, 0.18, AMBER)
txt(s, 0.9, 1.5, 11.5, 1.0, [("一句话总结", 24, AMBER, True)])
txt(s, 0.9, 2.4, 11.6, 2.5, [
    ("我们已经有了一条完整的、判别式的数字井筒——", 26, WHITE, True),
    ("9 条连续曲线、5 段统计、清晰的甜点与梯度。", 26, WHITE, True),
    ("", 12, WHITE, False),
    ("缺的不是数据的“量”，而是“生成”这一跳：", 24, ICE, False),
    ("把几块岩心，补成一整条可外推、带不确定性的虚拟井筒。", 24, ICE, False),
], space=1.2)
txt(s, 0.92, 6.6, 11.5, 0.5, [("GJ5-15 · 生成式数字井筒研究", 13, MUTE, False)])

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(">> 已生成", os.path.relpath(OUT, ROOT), "  共", len(prs.slides._sldIdLst), "页")
